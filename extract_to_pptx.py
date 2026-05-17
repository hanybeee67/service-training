import re
from pptx import Presentation
from pptx.util import Inches, Pt
from bs4 import BeautifulSoup

def clean_html(html):
    soup = BeautifulSoup(html, "html.parser")
    text = soup.get_text(separator="\n")
    # clean up multiple newlines
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()

with open('index.html', 'r', encoding='utf-8') as f:
    content = f.read()

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Everest Service Training Manual"
slide.placeholders[1].text = "Content for Editing"

def extract_and_add(pattern, title_prefix, content_str):
    matches = re.finditer(pattern, content_str, flags=re.DOTALL)
    for m in matches:
        if len(m.groups()) == 3:
            ch_id, name, html = m.groups()
            label = f"Chapter {ch_id}"
        elif len(m.groups()) == 4:
            ch_id, label, name, html = m.groups()
        
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"[{title_prefix}] {label}: {name}"
        tf = slide.placeholders[1].text_frame
        tf.text = clean_html(html)
        # make font smaller to fit
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)

# Chapters EN
en_section = content[:content.find('var CH_NE')]
extract_and_add(r'(\d+):\{label:"([^"]+)",name:"([^"]+)",html:function\(\)\{return `(.*?)`\}\}', 'EN', en_section)

# Chapters NE
ne_section = content[content.find('var CH_NE'):]
extract_and_add(r'(\d+):\{label:"([^"]+)",name:"([^"]+)",html:function\(\)\{return `(.*?)`\}\}', 'NE', ne_section)

# Screens EN
en_screens = re.search(r'var EN_SCREENS\s*=\s*\{(.*?)\};', content, flags=re.DOTALL)
if en_screens:
    for m in re.finditer(r'([a-zA-Z0-9_]+):\s*`(.*?)`', en_screens.group(1), flags=re.DOTALL):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"[EN Screen] {m.group(1)}"
        tf = slide.placeholders[1].text_frame
        tf.text = clean_html(m.group(2))
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)

# Screens NE
ne_screens = re.search(r'var NE_SCREENS\s*=\s*\{(.*?)\};', content, flags=re.DOTALL)
if ne_screens:
    for m in re.finditer(r'([a-zA-Z0-9_]+):\s*`(.*?)`', ne_screens.group(1), flags=re.DOTALL):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"[NE Screen] {m.group(1)}"
        tf = slide.placeholders[1].text_frame
        tf.text = clean_html(m.group(2))
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)

prs.save('training_manual.pptx')
print("Created training_manual.pptx")
